import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import re

# Function to create a PowerPoint slide from the provided data
def create_ppt_slide(data):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title of the slide
    title = slide.shapes.title
    title.text = "Crisis Scenario Overview"

    # Add a new slide for content overview
    content_slide_layout = prs.slide_layouts[5]  # Blank slide layout
    content_slide = prs.slides.add_slide(content_slide_layout)

    # Define the positions for each field
    positions = {
        "scenario": (Inches(0.5), Inches(0.5)),
        "inciting-incident": (Inches(0.5), Inches(1.5)),
        "training-objectives": (Inches(0.5), Inches(2.5)),
        "channels": (Inches(0.5), Inches(4.0)),
        "stakeholders": (Inches(0.5), Inches(5.5)),
        "roles": (Inches(5.0), Inches(0.5)),
        "actions": (Inches(5.0), Inches(2.5)),
        "research": (Inches(5.0), Inches(4.0)),
        "info": (Inches(5.0), Inches(5.5)),
        "guidance": (Inches(5.0), Inches(7.0))
    }

    width = Inches(4.5)
    height = Inches(1.0)

    for index, row in data.iterrows():
        if row['Field'] in positions:
            left, top = positions[row['Field']]

            # Add field name as bold heading
            title_box = content_slide.shapes.add_textbox(left, top, width, Inches(0.3))
            title_frame = title_box.text_frame
            title_frame.text = row['Field']
            title_frame.paragraphs[0].font.size = Pt(12)
            title_frame.paragraphs[0].font.bold = True

            # Clean details text and format bold parts
            details = row['Summary'].replace("_x000D_", "")
            details = re.sub(r"\*\*(.*?)\*\*", lambda match: match.group(1).upper(), details)

            # Add details in a text box below the heading
            text_box = content_slide.shapes.add_textbox(left, top + Inches(0.4), width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = details
            p.font.size = Pt(9)
    
    return prs

# Streamlit app to take CSV input and generate a PowerPoint slide
st.title("PowerPoint Slide Creator")

data_file = st.file_uploader("Upload CSV File", type=["csv"])

if data_file is not None:
    # Read the CSV file
    try:
        # Attempt to read the CSV file with more robust handling
        df = pd.read_csv(data_file, encoding='utf-8', skip_blank_lines=True, engine='python')
        st.write("### Uploaded Data:")
        st.write(df)

        # Check if the required columns are present
        if 'Field' in df.columns and 'Summary' in df.columns:
            # Button to generate PowerPoint
            if st.button("Create PowerPoint Slide"):
                ppt = create_ppt_slide(df[['Field', 'Summary']])
                
                # Save to BytesIO object
                ppt_io = BytesIO()
                ppt.save(ppt_io)
                ppt_io.seek(0)
                
                # Provide download link
                st.download_button(label="Download PowerPoint Slide", 
                                   data=ppt_io, 
                                   file_name="scenario_overview.pptx", 
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else:
            st.error("CSV must contain 'Field' and 'Summary' columns.")
    except pd.errors.ParserError:
        st.error("Error reading file: The CSV file appears to be badly formatted. Please check for missing quotes or incorrect delimiters.")
    except Exception as e:
        st.error(f"Error reading file: {e}")
